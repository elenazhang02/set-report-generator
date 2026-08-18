"""
SET Program Execution Report — PPT Generator
=====================================================
Generates a supplier-facing PowerPoint deck from a SET/LPM image CSV export.

USAGE
-----
1. Set CSV_PATH to your exported image CSV.
2. Set SUPPLIER_NAME — the portion of tracker_name that is the supplier
   (e.g. "CAMPARI AMERICA"). Everything after it becomes the program name.
   If left blank the full tracker_name is used as the program name.
3. Optionally set SUPPLIER_LOGO_PATH to a local PNG/JPG of the supplier logo.
   Leave blank to attempt an auto-fetch from the web.
4. Run:  python3 build_supplier_ppt.py
5. The deck is saved next to the CSV with a timestamped filename.
"""

VERSION = '1'

import csv
import os
import re
import datetime
import warnings
import io
import shutil
warnings.filterwarnings('ignore')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image as PILImage

# ══════════════════════════════════════════════════════════════════════════════
# ── USER CONFIG  (edit these) ─────────────────────────────────────────────────
# ══════════════════════════════════════════════════════════════════════════════

CSV_PATH           = ''
SUPPLIER_NAME      = 'CAMPARI AMERICA'
SUPPLIER_LOGO_PATH = ''
OUTPUT_DIR         = ''

# ══════════════════════════════════════════════════════════════════════════════
# ── CONSTANTS ─────────────────────────────────────────────────────────────────
# ══════════════════════════════════════════════════════════════════════════════
SCRIPT_DIR  = os.path.dirname(os.path.abspath(__file__))
LOGO_SG     = os.path.join(SCRIPT_DIR, 'logos', 'sg', 'sg_logo.png')
LOGOS_DIR   = os.path.join(SCRIPT_DIR, 'logos')

LOGO_MIN_PX = 80

# ── Palette (updated to match TEMPLATE.pptx) ─────────────────────────────────
GOLD        = RGBColor(0xD8, 0xC1, 0x86)   # warm light gold — rules, accents, labels
GOLD_STRIPE = RGBColor(0xAA, 0x8A, 0x3A)   # darker gold — title slide left stripe
GOLD_LABEL  = RGBColor(0x76, 0x5F, 0x26)   # brown-gold — KPI subtitles, table headers
CHARCOAL    = RGBColor(0x1C, 0x1A, 0x18)   # title slide bg, footer
HDR_DARK    = RGBColor(0x2A, 0x2A, 0x29)   # header bar on overview + image slides
PANEL_DARK  = RGBColor(0x32, 0x32, 0x31)   # metadata panel on image slides
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
CREAM       = RGBColor(0xFB, 0xF8, 0xF2)
OFF_WHITE   = RGBColor(0xF2, 0xEF, 0xE8)
MID_GREY    = RGBColor(0x6B, 0x65, 0x5A)
MUTED       = RGBColor(0x88, 0x80, 0x68)
DARK_RED    = RGBColor(0xA1, 0x15, 0x1A)
LIGHT_GOLD  = RGBColor(0xF7, 0xEE, 0xCC)

W = Inches(13.333)
H = Inches(7.5)


# ══════════════════════════════════════════════════════════════════════════════
# ── DATA EXTRACTION ───────────────────────────────────────────────────────────
# ══════════════════════════════════════════════════════════════════════════════

REQUIRED_COLUMNS = {'image_id', 'image_url', 'tracker_name', 'mission_completed_dt'}

def load_csv(path):
    with open(path, newline='', encoding='utf-8-sig') as f:
        rows = list(csv.DictReader(f))
    if not rows:
        raise ValueError('The CSV file is empty.')
    missing = REQUIRED_COLUMNS - set(rows[0].keys())
    if missing:
        raise ValueError(
            f'CSV is missing required column(s): {", ".join(sorted(missing))}.\n'
            f'Columns found: {", ".join(rows[0].keys())}'
        )
    return rows

def _get(row, key, default=''):
    return row.get(key, default).strip() if row.get(key) else default

def derive_metadata(rows, supplier_override):
    tracker_name = rows[0]['tracker_name'].strip()
    ref_no       = _get(rows[0], 'reference_no', 'N/A')

    if supplier_override:
        prefix = supplier_override.strip().upper()
        upper  = tracker_name.upper()
        if upper.startswith(prefix):
            supplier = supplier_override.strip().title()
            program  = tracker_name[len(prefix):].strip()
        else:
            supplier = supplier_override.strip().title()
            program  = tracker_name
    else:
        supplier = ''
        program  = tracker_name

    states = sorted(set(_get(r, 'site_state', '?') for r in rows))
    market = '  ·  '.join(s for s in states if s != '?') or 'N/A'

    dates = sorted(set(r['mission_completed_dt'][:10] for r in rows))
    if len(dates) == 1:
        d = datetime.datetime.strptime(dates[0], '%Y-%m-%d')
        period = d.strftime('%B %d, %Y')
        period_short = d.strftime('%b%Y')
    else:
        d0 = datetime.datetime.strptime(dates[0],  '%Y-%m-%d')
        d1 = datetime.datetime.strptime(dates[-1], '%Y-%m-%d')
        if d0.year == d1.year and d0.month == d1.month:
            period = f"{d0.strftime('%B %d')} – {d1.strftime('%d, %Y')}"
        else:
            period = f"{d0.strftime('%b %d')} – {d1.strftime('%b %d, %Y')}"
        period_short = d0.strftime('%b%Y')

    images   = len(rows)
    accounts = len(set(_get(r, 'customer_name', 'Unknown') for r in rows))
    cities   = len(set(_get(r, 'city', 'Unknown') for r in rows))
    days     = len(dates)

    from collections import Counter
    type_breakdown = Counter(_get(r, 'image_type', 'Unknown') for r in rows)
    def banner(name):
        return name.split('#')[0].strip().title().replace("'S", "'s")
    banner_breakdown = Counter(banner(_get(r, 'customer_name', 'Unknown')) for r in rows)

    city_counter = Counter(_get(r, 'city', 'Unknown').title() for r in rows)
    top7_cities = city_counter.most_common(7)
    others_city_count = sum(cnt for c, cnt in city_counter.items()
                            if c not in dict(top7_cities))
    city_breakdown = dict(top7_cities)
    if others_city_count:
        city_breakdown['Others'] = others_city_count
    city_list = list(dict(top7_cities).keys()) + (['Others'] if others_city_count else [])

    top7_banners = banner_breakdown.most_common(7)
    others_banner_count = sum(cnt for b, cnt in banner_breakdown.items()
                              if b not in dict(top7_banners))
    banner_breakdown_top7 = dict(top7_banners)
    if others_banner_count:
        banner_breakdown_top7['Others'] = others_banner_count

    return dict(
        tracker_name=tracker_name,
        supplier=supplier,
        program=program,
        ref_no=ref_no,
        market=market,
        period=period,
        period_short=period_short,
        images=images,
        accounts=accounts,
        cities=cities,
        days=days,
        type_breakdown=dict(type_breakdown.most_common()),
        banner_breakdown=banner_breakdown_top7,
        city_breakdown=city_breakdown,
        city_list=city_list,
        states=states,
    )

def detect_brands(rows):
    names = set()
    for row in rows:
        names.add(row['tracker_name'].strip())
    return sorted(names)[:5]


def build_filename(meta, generated_ts, output_dir, csv_path):
    def slug(s):
        return re.sub(r'[^A-Za-z0-9]+', '_', s).strip('_')
    supplier_slug = slug(meta['supplier']) if meta['supplier'] else slug(meta['program'][:20])
    program_slug  = slug(meta['program'])
    state_slug    = '_'.join(meta['states'])
    ts            = generated_ts.strftime('%Y%m%d_%H%M')
    ref_slug = slug(meta['ref_no']) if meta['ref_no'] and meta['ref_no'] != 'N/A' else ''
    parts = [p for p in [supplier_slug, program_slug, meta['period_short'], ref_slug, ts] if p]
    name = '_'.join(parts) + '.pptx'
    base = output_dir if output_dir else os.path.dirname(csv_path)
    return os.path.join(base, name)


# ══════════════════════════════════════════════════════════════════════════════
# ── LOGO FETCHING ─────────────────────────────────────────────────────────────
# ══════════════════════════════════════════════════════════════════════════════

def _name_slug(name):
    return re.sub(r'[^a-z0-9]+', '_', name.lower()).strip('_')


def _logo_usable(path):
    try:
        w, h = PILImage.open(path).size
        return min(w, h) >= LOGO_MIN_PX
    except Exception:
        return False


def find_local_logo(name):
    slug = _name_slug(name)
    search_dirs = [
        os.path.join(LOGOS_DIR, 'supplier'),
        os.path.join(LOGOS_DIR, 'brand'),
    ]
    for d in search_dirs:
        if not os.path.isdir(d):
            continue
        for fname in os.listdir(d):
            if not fname.lower().endswith('.png'):
                continue
            if _name_slug(os.path.splitext(fname)[0]) == slug:
                path = os.path.join(d, fname)
                if _logo_usable(path):
                    print(f'  Logo: local ({os.path.relpath(path, SCRIPT_DIR)}) ✓')
                    return path
                else:
                    print(f'  Logo: local file too small — falling back to web ({fname})')
    return None


def fetch_supplier_logo(supplier_name, cache_dir='/tmp/supplier_logos'):
    local = find_local_logo(supplier_name)
    if local:
        return local

    os.makedirs(cache_dir, exist_ok=True)
    safe = _name_slug(supplier_name)
    out  = os.path.join(cache_dir, f'{safe}.png')

    if os.path.exists(out):
        print(f'  Logo: web cache hit → {out}')
        return out

    try:
        import requests

        words   = supplier_name.lower().split()
        domains = [
            ''.join(words) + '.com',
            '-'.join(words) + '.com',
            words[0] + 'group.com',
            words[0] + '.com',
        ]
        for domain in domains:
            r = requests.get(f'https://logo.clearbit.com/{domain}', timeout=8)
            if r.status_code == 200 and 'image' in r.headers.get('Content-Type', ''):
                PILImage.open(io.BytesIO(r.content)).convert('RGBA').save(out, 'PNG')
                print(f'  Logo: Clearbit ({domain}) ✓')
                return out

        search_url = ('https://en.wikipedia.org/api/rest_v1/page/summary/' +
                      supplier_name.replace(' ', '_'))
        r = requests.get(search_url, timeout=8,
                         headers={'User-Agent': 'SGWSReportBot/1.0'})
        if r.status_code == 200:
            data  = r.json()
            thumb = (data.get('originalimage') or data.get('thumbnail') or {}).get('source')
            if thumb:
                r2 = requests.get(thumb, timeout=10)
                if r2.status_code == 200:
                    PILImage.open(io.BytesIO(r2.content)).convert('RGBA').save(out, 'PNG')
                    print(f'  Logo: Wikipedia thumbnail ✓')
                    return out

    except Exception as e:
        print(f'  Logo fetch error: {e}')

    print(f'  Logo: not found — will use text fallback for "{supplier_name}"')
    return None


# ══════════════════════════════════════════════════════════════════════════════
# ── PPT DRAWING HELPERS ───────────────────────────────────────────────────────
# ══════════════════════════════════════════════════════════════════════════════

def rect(slide, x, y, w, h, fill, line_color=None, line_w=0):
    sh = slide.shapes.add_shape(1, int(x), int(y), int(w), int(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = line_w
    else:
        sh.line.fill.background()
    return sh

def txt(slide, text, x, y, w, h, size, color,
        bold=False, italic=False, align=PP_ALIGN.LEFT, font='Calibri'):
    tb = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text           = text
    r.font.size      = Pt(size)
    r.font.color.rgb = color
    r.font.bold      = bold
    r.font.italic    = italic
    r.font.name      = font

def place_logo(slide, path, x, y, target_h, max_w=None):
    if not path or not os.path.exists(path):
        return 0
    img = PILImage.open(path)
    iw, ih = img.size
    w = target_h * (iw / ih)
    if max_w and w > max_w:
        w = max_w
        target_h = w * (ih / iw)
    slide.shapes.add_picture(path, int(x), int(y), width=int(w), height=int(target_h))
    return w

def img_fit(slide, path, x, y, bw, bh):
    im = PILImage.open(path)
    iw, ih = im.size
    scale = min(bw / iw, bh / ih)
    nw, nh = iw * scale, ih * scale
    slide.shapes.add_picture(path,
        int(x + (bw - nw) / 2), int(y + (bh - nh) / 2),
        width=int(nw), height=int(nh))

def gold_rule(slide, y, x=0, w=None, h=Inches(0.055)):
    rect(slide, x, y, w or W, h, GOLD)

def sg_logo(slide):
    """Small square SG logo, top-right corner."""
    place_logo(slide, LOGO_SG,
               W - Inches(0.858), Inches(0.217),
               Inches(0.558), max_w=Inches(0.558))

def footer(slide, page=None):
    fy, fh = H - Inches(0.36), Inches(0.36)
    rect(slide, 0, fy, W, fh, CHARCOAL)
    txt(slide, 'Southern Glazer\'s Wine & Spirits  ·  Confidential',
        Inches(0.3), fy + Inches(0.07), Inches(7), Inches(0.26),
        8, RGBColor(0xB8, 0xAD, 0x90))
    if page:
        txt(slide, str(page),
            W - Inches(0.55), fy + Inches(0.07), Inches(0.4), Inches(0.26),
            8, RGBColor(0xB8, 0xAD, 0x90), align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# ── SLIDE BUILDERS ────────────────────────────────────────────────────────────
# ══════════════════════════════════════════════════════════════════════════════

def build_title_slide(prs, meta, supplier_logo, ts_display):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)

    rect(s, 0, 0, W, H, CHARCOAL)
    rect(s, 0, 0, Inches(0.11), H, GOLD_STRIPE)              # left stripe (darker gold)

    HDR = Inches(1.1)
    rect(s, Inches(0.11), 0, W - Inches(0.11), HDR,
         RGBColor(0x14, 0x12, 0x10))
    gold_rule(s, HDR)

    # Supplier logo — text fallback if unavailable
    if supplier_logo and os.path.exists(supplier_logo):
        place_logo(s, supplier_logo, Inches(0.25), Inches(0.27),
                   Inches(0.56), max_w=Inches(2.4))
    else:
        fallback = (meta['supplier'] or meta['program']).upper()
        txt(s, fallback,
            Inches(0.25), Inches(0.3), Inches(4.5), Inches(0.5),
            15, GOLD, bold=True, font='Calibri')

    sg_logo(s)

    # Supplier + program names
    name_display = meta['supplier'] if meta['supplier'] else meta['program']
    txt(s, name_display.upper(),
        Inches(0.25), Inches(1.452), Inches(11), Inches(0.976),
        52, WHITE, bold=True, font='Calibri')
    txt(s, meta['program'],
        Inches(0.25), Inches(2.379), Inches(9), Inches(0.673),
        34, RGBColor(0xF0, 0xE0, 0xA8), italic=False, font='Calibri Light')

    gold_rule(s, Inches(3.32))

    meta_row = [
        ('SUPPLIER', meta['supplier'] or meta['program']),
        ('PROGRAM',  meta['program']),
        ('MARKET',   meta['market']),
        ('PERIOD',   meta['period']),
        ('REFERENCE NUMBER', meta['ref_no']),
    ]
    for i, (lbl, val) in enumerate(meta_row):
        cx = Inches(0.25) + i * Inches(2.6)
        txt(s, lbl, cx, Inches(3.607), Inches(2.5), Inches(0.303), 12, GOLD, bold=True)
        txt(s, val,  cx, Inches(3.847), Inches(2.5), Inches(0.303), 12, WHITE)

    txt(s, 'SET PROGRAM  ·  EXECUTION REPORT',
        Inches(0.25), Inches(4.367), Inches(9), Inches(0.303), 12, GOLD, bold=True)
    txt(s, f'Generated  {ts_display}',
        Inches(0.25), Inches(4.747), Inches(6), Inches(0.269), 10, GOLD)


def build_overview_slide(prs, meta, ts_display, brand_logos=None):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, W, H, CREAM)

    # Header — starts flush at top (no top gold rule)
    HDR = Inches(1.005)
    rect(s, 0, 0, W, HDR, HDR_DARK)
    sg_logo(s)

    program_display = (f'{meta["supplier"].upper()}  ·  {meta["program"]}'
                       if meta['supplier'] else meta['program'])
    txt(s, program_display,
        Inches(0.35), Inches(0.161), Inches(8.5), Inches(0.438),
        20, WHITE, bold=True)
    txt(s, f'Program Execution — Overview  ·  Generated {ts_display}',
        Inches(0.35), Inches(0.551), Inches(10), Inches(0.303), 12, GOLD)

    # KPI boxes
    kpis = [
        (str(meta['images']),   'IMAGES\nCAPTURED'),
        (str(meta['accounts']), 'RETAIL\nACCOUNTS'),
        (str(meta['cities']),   'CITIES\nCOVERED'),
        (str(meta['days']),     'DAYS OF\nEXECUTION'),
    ]
    bw, bh = Inches(2.8), Inches(1.62)
    by = Inches(1.309)
    gap = Inches(0.14)
    sx = (W - (4*bw + 3*gap)) / 2
    for i, (num, lbl) in enumerate(kpis):
        bx = sx + i*(bw+gap)
        rect(s, bx, by, bw, bh, WHITE,
             line_color=RGBColor(0xD8, 0xCC, 0xA8), line_w=Pt(1))
        rect(s, bx, by, bw, Inches(0.06), GOLD)
        txt(s, num, bx, by+Inches(0.148), bw, Inches(0.880),
            46, DARK_RED, bold=True, align=PP_ALIGN.CENTER)
        txt(s, lbl, bx, by+Inches(0.920), bw, Inches(0.471),
            11, GOLD_LABEL, bold=True, align=PP_ALIGN.CENTER)

    # Two breakdown tables
    def table(slide, title, data, x, y, w):
        txt(slide, title, x, y, w, Inches(0.269), 10, GOLD_LABEL, bold=True)
        rect(slide, x, y+Inches(0.260), w, Inches(0.025), GOLD)
        for i, (lbl, cnt) in enumerate(data):
            ry = y + Inches(0.350) + i*Inches(0.440)
            rect(slide, x, ry, w, Inches(0.380),
                 OFF_WHITE if i % 2 == 0 else WHITE)
            txt(slide, lbl, x+Inches(0.1), ry+Inches(0.06),
                w-Inches(0.7), Inches(0.286), 11, CHARCOAL)
            txt(slide, str(cnt), x+w-Inches(0.55), ry+Inches(0.06),
                Inches(0.45), Inches(0.286), 11, CHARCOAL,
                bold=True, align=PP_ALIGN.RIGHT)

    ty   = Inches(3.02)
    tw   = Inches(4.0)
    tgap = Inches(0.27)
    tx0  = Inches(0.4)
    table(s, 'IMAGE TYPE BREAKDOWN',
          list(meta['type_breakdown'].items()),
          tx0, ty, tw)
    table(s, 'RETAIL BANNER BREAKDOWN',
          list(meta['banner_breakdown'].items()),
          tx0 + tw + tgap, ty, tw)
    table(s, 'TOP CITIES',
          list(meta['city_breakdown'].items()),
          tx0 + 2 * (tw + tgap), ty, tw)

    # Brand logos strip
    if brand_logos:
        txt(s, 'FEATURED BRANDS',
            Inches(0.4), Inches(6.149), Inches(1.8), Inches(0.269), 10, GOLD_LABEL, bold=True)
        lx = Inches(2.1)
        for path in brand_logos.values():
            if path and os.path.exists(path):
                pw = place_logo(s, path, lx, Inches(6.129), Inches(0.554), max_w=Inches(2.0))
                lx += pw + Inches(0.32)

    footer(s, 2)


def build_image_slide(prs, row, img_path, page_num):
    blank = prs.slide_layouts[6]
    sl = prs.slides.add_slide(blank)
    rect(sl, 0, 0, W, H, CREAM)

    # Header — flush to top
    HDR_H = Inches(1.005)
    rect(sl, 0, 0, W, HDR_H, HDR_DARK)

    store = _get(row, 'customer_name', 'Unknown').title().replace("'S", "'s")
    city  = _get(row, 'city', '').title()
    state = _get(row, 'site_state', '')
    location = '  ·  '.join(filter(None, [store, f'{city}, {state}'.strip(', ')]))
    txt(sl, location,
        Inches(0.14), Inches(0.290), Inches(9.5), Inches(0.438),
        20, WHITE, bold=True)
    sg_logo(sl)
    gold_rule(sl, HDR_H)

    PANEL_W   = Inches(3.0)
    panel_top = Inches(0.990)
    panel_h   = H - panel_top - Inches(0.36)
    rect(sl, 0, panel_top, PANEL_W, panel_h, PANEL_DARK)

    # Build panel fields dynamically — only include rows whose column exists in the CSV
    cols = set(row.keys())
    candidate_items = [
        ('PROGRAM',          'tracker_name',        _get(row, 'tracker_name', '')),
        ('REFERENCE NUMBER', 'reference_no',         _get(row, 'reference_no', '')),
        ('DATE',             'mission_completed_dt', row['mission_completed_dt'][:10]),
        ('IMAGE TYPE',       'image_type',           _get(row, 'image_type', '')),
        ('SCENE TYPE',       'scene_type',           _get(row, 'scene_type', '')),
        ('MARKET',           'site_name',            '  ·  '.join(filter(None, [_get(row, 'site_name', ''), state]))),
    ]
    meta_items = [(lbl, val) for lbl, col, val in candidate_items if col in cols and val]

    # Distribute items evenly across available panel height
    panel_usable = panel_h - Inches(0.18) - Inches(0.10)
    item_h = panel_usable / max(len(meta_items), 1)
    lbl_size  = max(8,  min(10, int(10 * (item_h / Inches(0.65)))))
    val_size  = max(10, min(12, int(12 * (item_h / Inches(0.65)))))

    my = panel_top + Inches(0.18)
    for lbl, val in meta_items:
        txt(sl, lbl, Inches(0.14), my,
            PANEL_W-Inches(0.2), Inches(0.25), lbl_size, GOLD, bold=True)
        txt(sl, val, Inches(0.14), my + Inches(0.21),
            PANEL_W-Inches(0.2), item_h - Inches(0.25), val_size, WHITE)
        my += item_h

    img_x = PANEL_W + Inches(0.06)
    img_y = panel_top
    img_w = W - img_x
    img_h = panel_h
    img_fit(sl, img_path, img_x+Inches(0.1), img_y+Inches(0.1),
            img_w-Inches(0.2), img_h-Inches(0.2))

    footer(sl, page_num)


# ══════════════════════════════════════════════════════════════════════════════
# ── IMAGE DOWNLOADING ─────────────────────────────────════════════════════════
# ══════════════════════════════════════════════════════════════════════════════

def download_images(rows, csv_path, workers=32):
    from concurrent.futures import ThreadPoolExecutor, as_completed
    import threading

    cache_dir = os.path.join(os.path.dirname(os.path.abspath(csv_path)), '.img_cache')
    os.makedirs(cache_dir, exist_ok=True)

    import glob
    sentinel = os.path.join(cache_dir, '.last_cleared')
    today = datetime.date.today().isoformat()
    last_cleared = open(sentinel).read().strip() if os.path.exists(sentinel) else ''
    if last_cleared != today:
        for f in glob.glob(os.path.join(cache_dir, '*.jpg')):
            os.remove(f)
        open(sentinel, 'w').write(today)
        print(f'Cache cleared (first run today)')

    try:
        import requests
    except ImportError:
        print('ERROR: requests not installed — run: pip3 install requests')
        return {}

    total = len(rows)
    print(f'Images  : {total} total  |  cache → {cache_dir}  |  workers → {workers}')

    img_files  = {}
    lock       = threading.Lock()
    counters   = {'downloaded': 0, 'cached': 0, 'failed': 0}

    def fetch(args):
        i, row = args
        img_id = row['image_id']
        path   = os.path.join(cache_dir, f'{img_id}.jpg')

        if os.path.exists(path):
            with lock:
                img_files[img_id] = path
                counters['cached'] += 1
                print(f'  [{i+1:>{len(str(total))}}/{total}] cache   {row["customer_name"]} ({row["city"]})')
            return

        try:
            r = requests.get(row['image_url'], timeout=25)
            if r.status_code == 200:
                img = PILImage.open(io.BytesIO(r.content)).convert('RGB')
                img.save(path, 'JPEG', quality=85)
                with lock:
                    img_files[img_id] = path
                    counters['downloaded'] += 1
                    print(f'  [{i+1:>{len(str(total))}}/{total}] fetched {row["customer_name"]} ({row["city"]})  {len(r.content)//1024}KB')
            else:
                with lock:
                    counters['failed'] += 1
                    print(f'  [{i+1:>{len(str(total))}}/{total}] HTTP {r.status_code} — {row["customer_name"]}')
        except Exception as e:
            with lock:
                counters['failed'] += 1
                print(f'  [{i+1:>{len(str(total))}}/{total}] ERROR — {row["customer_name"]}: {e}')

    with ThreadPoolExecutor(max_workers=workers) as pool:
        futures = [pool.submit(fetch, (i, row)) for i, row in enumerate(rows)]
        for f in as_completed(futures):
            f.result()

    print(f'  → {counters["downloaded"]} fetched  |  {counters["cached"]} cached  |  {counters["failed"]} failed\n')
    return img_files


# ══════════════════════════════════════════════════════════════════════════════
# ── MAIN ──────────────────────────────────────────────────────────────────────
# ══════════════════════════════════════════════════════════════════════════════

def prompt_yes_no(question, default='n'):
    hint = '[Y/n]' if default == 'y' else '[y/N]'
    answer = input(f'{question} {hint}: ').strip().lower()
    if not answer:
        return default == 'y'
    return answer in ('y', 'yes')


def resolve_csv_paths(pick_first=False):
    if CSV_PATH and os.path.exists(CSV_PATH):
        return [CSV_PATH]

    script_dir = os.path.dirname(os.path.abspath(__file__))
    candidates = sorted(
        f for f in os.listdir(script_dir)
        if f.startswith('lpm-images-') and f.endswith('.csv')
    )

    if not candidates:
        raise FileNotFoundError(
            f'No lpm-images-*.csv found in {script_dir}\n'
            'Set CSV_PATH at the top of the script or drop the file in that folder.'
        )

    if len(candidates) == 1 or pick_first:
        path = os.path.join(script_dir, candidates[0])
        print(f'CSV   : {candidates[0]}  (auto-detected)')
        return [path]

    print(f'{len(candidates)} CSVs found — which would you like to process?\n')
    for i, name in enumerate(candidates, 1):
        print(f'  [{i}] {name}')
    print(f'  [A] All ({len(candidates)} decks)')
    print()
    while True:
        choice = input(f'Enter number (1–{len(candidates)}) or A for all: ').strip().lower()
        if choice == 'a':
            return [os.path.join(script_dir, f) for f in candidates]
        if choice.isdigit() and 1 <= int(choice) <= len(candidates):
            return [os.path.join(script_dir, candidates[int(choice) - 1])]
        print(f'  Please enter a number between 1 and {len(candidates)}, or A.')


def build_deck(csv_path, include_supplier_logo, include_brand_logos, generated_at):
    ts_display = generated_at.strftime('%B %d, %Y  %H:%M %Z')

    rows = load_csv(csv_path)
    print(f'\nCSV   : {os.path.basename(csv_path)}')
    print(f'Rows  : {len(rows)}')

    meta = derive_metadata(rows, SUPPLIER_NAME)
    print(f'Supplier : {meta["supplier"] or "(from tracker_name)"}')
    print(f'Program  : {meta["program"]}')
    print(f'Market   : {meta["market"]}')
    print(f'Period   : {meta["period"]}')
    print(f'Ref #    : {meta["ref_no"]}')

    supplier_logo = None
    if include_supplier_logo:
        if SUPPLIER_LOGO_PATH and os.path.exists(SUPPLIER_LOGO_PATH):
            supplier_logo = SUPPLIER_LOGO_PATH
            print(f'  Supplier logo: using local file')
        else:
            print(f'  Fetching supplier logo for "{meta["supplier"] or meta["program"]}"...')
            supplier_logo = fetch_supplier_logo(meta['supplier'] or meta['program'])
            if not supplier_logo:
                print('  → Not found — supplier name shown as text instead')
    else:
        print('  Supplier logo: skipped — name will be shown as text')

    brand_logos = {}
    if include_brand_logos:
        print('  Fetching brand logos...')
        brand_names = detect_brands(rows)
        for brand in brand_names:
            print(f'    {brand}...')
            path = fetch_supplier_logo(brand)
            if path:
                brand_logos[brand] = path
                print(f'    ✓ {brand}')
            else:
                print(f'    ✗ {brand} — skipped')
    else:
        print('  Brand logos: skipped')

    print()
    img_files = download_images(rows, csv_path)

    print('Building slides...')
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    build_title_slide(prs, meta, supplier_logo, ts_display)
    build_overview_slide(prs, meta, ts_display, brand_logos)

    for idx, row in enumerate(rows):
        img_id   = row['image_id']
        img_path = img_files.get(img_id)
        if img_path:
            build_image_slide(prs, row, img_path, idx + 3)
        else:
            print(f'  SKIP slide {idx+3} — no image for {img_id}')

    out_path = build_filename(meta, generated_at, OUTPUT_DIR, csv_path)
    prs.save(out_path)

    print(f'\n✓  Saved → {out_path}')
    print(f'   Slides : {len(prs.slides)}')
    print(f'   Size   : {os.path.getsize(out_path) // (1024*1024)} MB')

    cache_dir = os.path.join(os.path.dirname(os.path.abspath(csv_path)), '.img_cache')
    if os.path.isdir(cache_dir):
        shutil.rmtree(cache_dir)
        print(f'   Cache  : cleared ({cache_dir})')

    return out_path


def main():
    import sys
    madmax = '--madmax' in sys.argv

    # --csv PATH overrides auto-resolution
    explicit_csv = None
    if '--csv' in sys.argv:
        idx = sys.argv.index('--csv')
        if idx + 1 < len(sys.argv):
            explicit_csv = sys.argv[idx + 1]

    generated_at = datetime.datetime.now().astimezone()

    print(f'\n{"="*60}')
    print(f'SET Program Execution Report — Generator  [v{VERSION}]')
    if madmax:
        print('⚡ MAD MAX MODE — no prompts, first CSV, no logos')
    print(f'{"="*60}')

    if explicit_csv:
        csv_paths = [explicit_csv]
    elif madmax:
        csv_paths = resolve_csv_paths(pick_first=True)
    else:
        csv_paths = resolve_csv_paths()

    if madmax:
        include_supplier_logo = False
        include_brand_logos   = False
    else:
        print()
        include_supplier_logo = prompt_yes_no(
            'Include supplier logo on title slide?\n'
            '  (pulled from web if not cached — adds ~5-15s)',
            default='y'
        )
        include_brand_logos = prompt_yes_no(
            'Include brand logos on overview slide?\n'
            '  (each logo pulled from web if not cached — adds ~10-30s)',
            default='n'
        )

    saved = []
    for csv_path in csv_paths:
        out = build_deck(csv_path, include_supplier_logo, include_brand_logos, generated_at)
        saved.append(out)

    if len(saved) > 1:
        print(f'\n{"="*60}')
        print(f'All done — {len(saved)} decks saved:')
        for p in saved:
            print(f'  {p}')
        print()

    return saved


if __name__ == '__main__':
    main()
